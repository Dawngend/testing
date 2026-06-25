import os
import io
import logging
import hashlib
import pdfplumber
import pytesseract
import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Suppress verbose pdfminer FontBBox warnings
logging.getLogger("pdfminer").setLevel(logging.ERROR)

# Tesseract executable path for Windows
#pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def get_cache_filename(file_path):
    """Creates a matching .txt filename next to the original file."""
    return file_path + "_saved_text.txt"

def extract_text_from_pdf(file_path, cache_file):
    extracted_text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = ""
                text = page.extract_text()
                if text:
                    page_text += text + "\n\n"

                # Scan full-page screenshots
                if not text or len(text.strip()) < 50:
                    print(f"  -> Scanning image on PDF page {i + 1}...")
                    try:
                        pil_image = page.to_image(resolution=300).original
                        ocr_text = pytesseract.image_to_string(pil_image)
                        if ocr_text.strip():
                            page_text += "\n[Scanned from Image]:\n" + ocr_text + "\n\n"
                    except Exception as e:
                        print(f"  [Warning] Failed to OCR image on page {i + 1}: {e}")
                
                extracted_text += page_text
                
                # INCREMENTAL SAVE: Write to hard drive immediately!
                with open(cache_file, "a", encoding="utf-8") as f:
                    f.write(page_text)

        return extracted_text
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_pptx(file_path, cache_file):
    extracted_text = ""
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
                
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f"  -> Scanning image on PPTX slide {i + 1}...")
                    try:
                        image_bytes = shape.image.blob
                        image = Image.open(io.BytesIO(image_bytes))
                        ocr_text = pytesseract.image_to_string(image)
                        if ocr_text.strip():
                            slide_text += "\n[Scanned from Image]:\n" + ocr_text + "\n"
                    except Exception as e:
                        print(f"  [Warning] Failed to OCR PPTX image: {e}")
                        
            slide_text += "\n--- Next Slide ---\n\n"
            extracted_text += slide_text
            
            # INCREMENTAL SAVE: Write to hard drive immediately!
            with open(cache_file, "a", encoding="utf-8") as f:
                f.write(slide_text)

        return extracted_text
    except Exception as e:
        return f"Error reading PPTX: {e}"

@st.cache_data
def process_module_file_v2(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()
    
    cache_file = get_cache_filename(file_path)
    
    # Check if a saved text file already exists!
    if os.path.exists(cache_file):
        print(f"  [Fast Load] Found saved extraction! Skipping Tesseract for {file_path}...")
        with open(cache_file, "r", encoding="utf-8") as f:
            return f.read()

    # If no cache exists, make sure we start with a clean text file
    if os.path.exists(cache_file):
        os.remove(cache_file)

    if file_extension == '.pdf':
        print(f"Processing PDF: {file_path}...")
        return extract_text_from_pdf(file_path, cache_file)
    elif file_extension == '.pptx':
        print(f"Processing PPTX: {file_path}...")
        return extract_text_from_pptx(file_path, cache_file)
    else:
        return "Unsupported file type. Please upload a PDF or PPTX."

def process_file(file_path: str, user_id: str) -> dict:
    """Processes a PDF/PPTX file, extracts text, and returns text and metadata with file hash."""
    text = process_module_file_v2(file_path)
    
    # Calculate hash of file to use as file_hash/doc_id
    sha256 = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256.update(byte_block)
        file_hash = sha256.hexdigest()
    except Exception as e:
        print(f"  [Warning] Failed to calculate hash for {file_path}: {e}")
        # Fallback to hash of file_path if file read fails
        file_hash = hashlib.sha256(file_path.encode()).hexdigest()
        
    return {
        'text': text,
        'metadata': {
            'file_hash': file_hash,
            'user_id': user_id
        }
    }

def chunk_text_for_rag(text: str, chunk_size: int = 1000, overlap: int = 200) -> list:
    """Split text into chunks of roughly chunk_size characters with some overlap."""
    if not text:
        return []
        
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        if end >= text_len:
            chunks.append(text[start:])
            break
            
        # Try to break at a space or newline to avoid cutting words
        last_space = text.rfind(' ', start, end)
        if last_space > start + (chunk_size - overlap):
            end = last_space
            
        chunks.append(text[start:end])
        start = end - overlap
        if start < 0:
            start = 0
            
    return chunks